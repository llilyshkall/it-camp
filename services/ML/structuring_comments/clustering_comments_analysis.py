import os
import docx
import fitz
import torch
import pandas as pd
import asyncio
import aiohttp
from tqdm.asyncio import tqdm as async_tqdm
import json
import subprocess
import shutil
from enum import StrEnum
from pathlib import Path
from typing import NamedTuple

import jinja2

from pptx import Presentation
from bs4 import BeautifulSoup

from langchain.schema.document import Document
from langchain_community.retrievers import BM25Retriever
from langchain.retrievers import EnsembleRetriever
from langchain_community.vectorstores import FAISS
from langchain_huggingface import HuggingFaceEmbeddings
from langchain.text_splitter import RecursiveCharacterTextSplitter

# ---  ГЛОБАЛЬНАЯ КОНФИГУРАЦИЯ  ---
GLOBAL_CONFIG = {
    "LOCAL_API_URL": "http://89.108.116.240:11434/api/chat",
    "LOCAL_MODEL_NAME": "qwen3:8b",
    "EMBEDDING_MODEL": "intfloat/multilingual-e5-large",
    "MAX_CONCURRENT_REQUESTS": 1,
    "RETRIEVER_TOP_K": 5,
    "REQUEST_DELAY_SECONDS": 0.5
}


# --- PDF generator ----
def filter_newlines(text: str) -> str:
    """Фильтрует переносы строк для LaTeX"""
    # Заменяем переносы строк на пробелы и удаляем лишние пробелы
    return ' '.join(text.replace("|", "").split()).replace('%', '\\%').replace('&', '\\&').replace("_", "\\_")


def process_summary(text: str, section_label: str) -> str:
    """Обрабатывает сводку, добавляя ссылки на строки таблицы"""
    # Простая замена номеров на ссылки (можно адаптировать под ваш формат)
    import re
    # Ищем упоминания источников типа "ИСТОЧНИК 1", "источник 3" и т.д.
    pattern = r'(ИСТОЧНИК|источник|Источник|Source|SOURCE|КОНТЕКСТ)\s+(\d+)'

    def replace_with_ref(match):
        source_num = match.group(2)
        return f"{match.group(1)} \\hyperlink{{row:{section_label}-{source_num}}}{{{source_num}}}"

    return re.sub(pattern, replace_with_ref, text)

jinja_env = jinja2.Environment(
    block_start_string='<BLOCK>',
    block_end_string='</BLOCK>',
    variable_start_string='<VAR>',
    variable_end_string='</VAR>',
    comment_start_string='<!--',
    comment_end_string='-->'
)

jinja_env.filters['filter_newlines'] = filter_newlines
jinja_env.filters['process_summary'] = process_summary

class Source(NamedTuple):
    source_label: str
    source_chunk: str
    source_filepath: str


class Section(NamedTuple):
    name: str
    label: str
    summary: str
    sources: list[Source]


class Status(StrEnum):
    confirmed = "confirmed"
    not_found = "not_found"
    partial = "partial"
    indirect = "indirect"
    requires_confirmation = "requires_confirmation"

    def get_ru_name(self) -> str:
        match self:
            case Status.confirmed:
                return "Подтверждено"
            case Status.not_found:
                return "Не найдено"
            case Status.partial:
                return "Частично найдено"
            case Status.indirect:
                return "Индиректно"
            case Status.requires_confirmation:
                return "Требует подтверждения"
            case _:
                raise NotImplementedError


CHAPTER_TEMPLATE = jinja_env.from_string("""
\\chapter{<VAR>chapter_name</VAR>}
\\label{cha:<VAR>chapter_label</VAR>}

<BLOCK>for section in sections</BLOCK>
\\section{<VAR>section.name</VAR>}
\\label{sec:<VAR>section.label</VAR>}

\\subsection{Краткая сводка}
\\textbf{Статус}: <VAR>chapter_name</VAR>. <VAR>section.summary | filter_newlines | process_summary(section.label)</VAR>

\\subsection{Релевантные результаты}
\\begin{longtable}{|p{0.05\\textwidth}|p{0.65\\textwidth}|p{0.2\\textwidth}|}
\\hline
\\textbf{№} & \\textbf{Релевантные результаты} & \\textbf{Источник} \\\\
\\hline
\\endhead
<BLOCK>for source in section.sources</BLOCK>
<BLOCK>set source_index = loop.index</BLOCK>
\\raisebox{-\\baselineskip}[0pt][0pt]{\\hypertarget{row:<VAR>section.label</VAR>-<VAR>source_index</VAR>}{}} <VAR>source_index</VAR> & <VAR>source.source_chunk | filter_newlines</VAR> & \\cite{<VAR>source.source_label</VAR>} \\\\
\\hline
<BLOCK>endfor</BLOCK>
\\end{longtable}
<BLOCK>endfor</BLOCK>
""")

BIBLIOGRAPHY_TEMPLATE = jinja_env.from_string("""@misc{<VAR>source_label</VAR>,
    author = {ПАО ``Газпром``},
    title = {<VAR>filename | filter_newlines</VAR>},
    howpublished = {Внутренний документ},
    year = {2025},
    note = {Страница: <VAR>page</VAR>, Слайд: <VAR>slide</VAR>}
}
""")

LATEX_ROOT_PATH = Path(__file__).parent / "latex-gost-template"
LATEX_OUTPUT_PATH = LATEX_ROOT_PATH / "thesis.pdf"
LATEX_COMPILE_SCRIPT = LATEX_ROOT_PATH / "build.sh"
LATEX_TEMPLATE_PATH = LATEX_ROOT_PATH / "tex"
LATEX_SOURCE_PATH = LATEX_ROOT_PATH / "tex_tmp"


def create_safe_label(text: str) -> str:
    """Создает безопасный лабел для LaTeX из текста"""
    return text.lower().replace(' ', '-').replace(',', '').replace('.', '').replace('(', '').replace(')', '')


def create_source_label(filename: str, index: int) -> str:
    """Создает уникальный идентификатор для источника"""
    base_name = Path(filename).stem.lower().replace(' ', '-').replace('.', '-')
    return f"{base_name}-{index}"


def render_latex(file_path: Path, **kwargs) -> None:
    template = jinja_env.from_string(file_path.read_text())
    rendered = template.render(**kwargs)
    file_path.write_text(rendered)


def make_pdf(input_json: dict, output_pdf: Path) -> None:
    # Создаем директории если они не существуют
    LATEX_SOURCE_PATH.mkdir(exist_ok=True)

    chapters = {status: [] for status in Status}
    json_data = input_json

    # Собираем все источники для библиографии
    all_sources = []
    source_counter = {}

    for section_name, content in json_data.items():
        status = Status(content["status"])

        # Обрабатываем источники для этой секции
        section_sources = []
        for i, source_data in enumerate(content["sources"]):
            filename = source_data["filename"]
            source_counter[filename] = source_counter.get(filename, 0) + 1
            source_label = create_source_label(filename, source_counter[filename])

            section_sources.append(Source(
                source_label=source_label,
                source_chunk=source_data["snippet"],
                source_filepath=filename
            ))

            # Добавляем в общий список источников
            all_sources.append({
                'source_label': source_label,
                'filename': filename,
                'page': source_data.get("page", "N/A") or "N/A",
                'slide': source_data.get("slide", "N/A") or "N/A"
            })

        # Создаем секцию
        section_label = create_safe_label(section_name)
        section = Section(
            name=section_name,
            label=section_label,
            summary=content["answer"],
            sources=section_sources
        )

        chapters[status].append(section)

    shutil.rmtree(LATEX_SOURCE_PATH)
    shutil.copytree(LATEX_TEMPLATE_PATH, LATEX_SOURCE_PATH)

    # Создаем библиографию
    bib_content = []
    for source in all_sources:
        bib_entry = BIBLIOGRAPHY_TEMPLATE.render(**source)
        bib_content.append(bib_entry)

    bib_path = LATEX_SOURCE_PATH / "0-main.bib"
    bib_path.write_text('\n'.join(bib_content), encoding='utf-8')

    # Создаем главы
    chapter_ids = []
    for i, status in enumerate(Status, start=1):
        if not chapters[status]:  # Пропускаем пустые главы
            continue

        chapter_id = f"chapter-{i}"
        chapter_filename = f"3{i}-{chapter_id}.tex"
        chapter_path = LATEX_SOURCE_PATH / chapter_filename

        rendered_template = CHAPTER_TEMPLATE.render(
            chapter_name=status.get_ru_name(),
            chapter_label=f"chapter-{i}",
            sections=chapters[status]
        )

        chapter_path.write_text(rendered_template, encoding='utf-8')
        chapter_ids.append(chapter_filename.replace('.tex', ''))

    render_latex(LATEX_SOURCE_PATH / "0-main.tex", chapters=chapter_ids)
    render_latex(LATEX_SOURCE_PATH / "11-title-page.tex", doc_name="Сводный отчёт по замечаниям и программе доизучения",
                 doc_id=" GAZ-REP-2025-001", assigned_to="Отдел аналитики")
    render_latex(LATEX_SOURCE_PATH / "2-intro.tex", intro_text="""Настоящий отчёт подготовлен на основании предоставленных
данных в формате JSON. Категории данных сформированы как
главы, группы — как подразделы. Для каждого подраздела
приведены синтезированное описание и исходные замечания""")
    render_latex(LATEX_SOURCE_PATH / "4-conclusion.tex", conclusion_text="""Предложенные мероприятия направлены на снижение
неопределённостей и повышение качества прогнозов.
Рекомендовано согласовать план доизучения и
актуализировать модели по итогам получения новых данных.""")

    # Компилируем LaTeX
    try:
        subprocess.run(['bash', str(LATEX_COMPILE_SCRIPT)], cwd=LATEX_SOURCE_PATH, check=True)

    except subprocess.CalledProcessError as e:
        print(f"Ошибка компиляции LaTeX: {e}")

    if LATEX_OUTPUT_PATH.exists():
        shutil.copy2(LATEX_OUTPUT_PATH, output_pdf)
        print(f"PDF успешно создан: {output_pdf}")
    else:
        print("Ошибка: PDF файл не был создан")

# --- RAG ---

class ComprehensiveRAGSystem:
    def __init__(self, config):
        self.config = config
        self.device = "cuda" if torch.cuda.is_available() else "cpu"
        print(f"Используется устройство: {self.device} для проекта '{self.config.get('PROJECT_NAME', 'Unknown')}'")

        print("Загрузка Embedding модели...")
        self.embedding_model = HuggingFaceEmbeddings(
            model_name=self.config["EMBEDDING_MODEL"],
            model_kwargs={'device': self.device},
            encode_kwargs={'normalize_embeddings': True}
        )
        print("✅ Embedding модель загружена.")

        self.retriever = self._build_or_load_retriever()

    def _extract_text_from_pptx_safe(self, filepath, filename):
        docs = []
        try:
            prs = Presentation(filepath)
            for i, slide in enumerate(prs.slides):
                slide_texts = [
                    shape.text for shape in slide.shapes
                    if hasattr(shape, "text") and shape.text
                ]
                if slide_texts:
                    slide_content = "\n".join(slide_texts)
                    docs.append(Document(
                        page_content=slide_content,
                        metadata={"filename": filename, "slide": i + 1}
                    ))
            return docs
        except Exception as e:
            print(f"    ⚠ PPTX (SVG или неподдерживаемый объект) – fallback через PyMuPDF: {e}")
            try:
                with fitz.open(filepath) as ppt_as_pdf:
                    for page_num, page in enumerate(ppt_as_pdf):
                        page_text = page.get_text("text")
                        if page_text:
                            docs.append(Document(
                                page_content=page_text,
                                metadata={"filename": filename, "page": page_num + 1}
                            ))
            except Exception as inner_e:
                print(f"    Не удалось fallback-распарсить {filename} через PyMuPDF: {inner_e}")
            return docs

    def _extract_text_from_docs(self, folder_path):
        if not os.path.exists(folder_path):
            print(f"    Папка с документами не найдена: {folder_path}")
            return []
        all_docs = []
        print(f"--- Начало извлечения текста (продвинутый парсер) из '{folder_path}' ---")

        for filename in os.listdir(folder_path):
            full_path = os.path.join(folder_path, filename)
            try:
                if filename.endswith(".docx"):
                    doc = docx.Document(full_path)
                    content_parts = []
                    for para in doc.paragraphs:
                        if para.text.strip(): content_parts.append(para.text)
                    for table in doc.tables:
                        for row in table.rows:
                            row_text = " | ".join([cell.text.strip() for cell in row.cells])
                            if row_text: content_parts.append(row_text)
                    full_text = "\n".join(content_parts)
                    if full_text: all_docs.append(Document(page_content=full_text, metadata={"filename": filename}))

                elif filename.endswith(".pdf"):
                    with fitz.open(full_path) as pdf_doc:
                        for page_num, page in enumerate(pdf_doc):
                            page_text = page.get_text("text")
                            if page_text: all_docs.append(
                                Document(page_content=page_text, metadata={"filename": filename, "page": page_num + 1}))

                elif filename.endswith(".pptx"):
                    docs_from_pptx = self._extract_text_from_pptx_safe(full_path, filename)
                    all_docs.extend(docs_from_pptx)

                elif filename.endswith((".html", ".htm")):
                    with open(full_path, 'r', encoding='utf-8', errors='ignore') as f:
                        html = f.read()
                    soup = BeautifulSoup(html, "lxml")
                    for tag in soup(["script", "style"]): tag.decompose()
                    text = soup.get_text(" ", strip=True)
                    if text: all_docs.append(Document(page_content=text, metadata={"filename": filename}))

                elif filename.endswith(".txt"):
                    with open(full_path, 'r', encoding='utf-8', errors='ignore') as f:
                        text = f.read()
                    if text: all_docs.append(Document(page_content=text, metadata={"filename": filename}))

                print(f"     Успешно обработан файл: {filename}")
            except Exception as e:
                print(f"     Не удалось прочитать файл {filename}: {e}")
        return all_docs

    def _split_text_into_chunks(self, documents):
        text_splitter = RecursiveCharacterTextSplitter(chunk_size=700, chunk_overlap=150)
        return text_splitter.split_documents(documents)

    def _build_or_load_retriever(self):
        index_path, chunks_path = self.config["INDEX_FILE_PATH"], self.config["CHUNKS_FILE_PATH"]
        index_dir = os.path.dirname(index_path)
        index_name = os.path.basename(index_path).replace('.faiss', '')

        if os.path.exists(index_path):
            print(f"--- Найден существующий индекс для '{self.config['PROJECT_NAME']}'. Загружаем... ---")
            faiss_store = FAISS.load_local(
                folder_path=index_dir, embeddings=self.embedding_model, index_name=index_name,
                allow_dangerous_deserialization=True
            )
            with open(chunks_path, 'r', encoding='utf-8') as f:
                chunks_json = json.load(f)
            split_docs = [Document(page_content=c["page_content"], metadata=c["metadata"]) for c in chunks_json]
        else:
            print(f"--- Индекс для '{self.config['PROJECT_NAME']}' не найден. Создаем новый... ---")
            documents = self._extract_text_from_docs(self.config["DOCUMENTS_PATH"])
            if not documents: raise FileNotFoundError(
                f"Документы для проекта '{self.config['PROJECT_NAME']}' не найдены в {self.config['DOCUMENTS_PATH']}.")
            split_docs = self._split_text_into_chunks(documents)
            print(f"\n--- Векторизуем {len(split_docs)} чанков для проекта '{self.config['PROJECT_NAME']}'... ---")
            faiss_store = FAISS.from_documents(split_docs, self.embedding_model)
            faiss_store.save_local(folder_path=index_dir, index_name=index_name)
            chunks_for_json = [{"page_content": doc.page_content, "metadata": doc.metadata} for doc in split_docs]
            with open(chunks_path, 'w', encoding='utf-8') as f:
                json.dump(chunks_for_json, f, ensure_ascii=False, indent=2)

        faiss_retriever = faiss_store.as_retriever(search_kwargs={"k": self.config["RETRIEVER_TOP_K"]})
        bm25_retriever = BM25Retriever.from_documents(split_docs)
        bm25_retriever.k = self.config["RETRIEVER_TOP_K"]

        ensemble_retriever = EnsembleRetriever(
            retrievers=[bm25_retriever, faiss_retriever], weights=[0.5, 0.5]
        )
        print(f"      Гибридный ретривер (BM25 + FAISS) готов для проекта '{self.config['PROJECT_NAME']}'.")
        return ensemble_retriever

    # --- ИЗМЕНЕНО: Добавлен параметр "format": "json" для принудительного вывода в JSON ---
    async def _call_local_llm(self, messages):
        payload = {
            "model": self.config["LOCAL_MODEL_NAME"],
            "messages": messages,
            "stream": False,
            "format": "json",
            "options": {
                "temperature": 0.2,  # Низкая температура для точности
                "top_p": 0.9,  # Оптимальная  выборка
                "repetition_penalty": 1.05  #  штраф за повторы
            }
        }
        try:
            async with aiohttp.ClientSession() as session:
                async with session.post(self.config["LOCAL_API_URL"], json=payload, timeout=300) as response:
                    if response.status != 200:
                        error_text = await response.text()
                        print(f"  [LLM SERVER ERROR] Статус {response.status}. Ответ: {error_text[:200]}")
                        # Возвращаем JSON с ошибкой, чтобы основной код мог это обработать
                        return json.dumps({"status": "requires_confirmation", "answer": f"Ошибка API: {error_text}"})
                    response_data = await response.json()
                    # Ответ от Ollama в JSON mode находится в message.content в виде строки
                    return response_data.get('message', {}).get('content', '')
        except Exception as e:
            print(f"  [LLM CONNECTION EXCEPTION] Ошибка при запросе к API: {type(e).__name__}: {e}")
            return json.dumps({"status": "requires_confirmation", "answer": f"Ошибка подключения к API: {e}"})

    async def _expand_query(self, criterion):
        prompt = f"""Переформулируй следующий запрос тремя разными способами для поиска в базе знаний. Используй синонимы и меняй структуру. Верни только 3 новые версии, каждая на новой строке НА РУССКОМ ЯЗЫКЕ.

ИСХОДНЫЙ ЗАПРОС: "{criterion}"

ПЕРЕФОРМУЛИРОВАННЫЕ ЗАПРОСЫ:"""
        response_text = await self._call_local_llm([{"role": "user", "content": prompt}])
        if "Ошибка" in response_text: return [criterion]
        expanded = [q.strip().lstrip("-* ").strip() for q in response_text.split('\n') if q.strip()]
        print(f"      ... запрос '{criterion[:30]}...' расширен до: {expanded}")
        return [criterion] + expanded



    async def process_criterion(self, criterion):
        print(f"\n--- Обработка критерия для '{self.config['PROJECT_NAME']}': '{criterion[:70]}...' ---")

        all_queries = await self._expand_query(criterion)
        tasks = [asyncio.to_thread(self.retriever.invoke, q) for q in all_queries]
        results_from_queries = await asyncio.gather(*tasks)

        unique_docs = {doc.page_content: doc for doc_list in results_from_queries for doc in doc_list}
        retrieved_docs = list(unique_docs.values())
        print(f"    🔍 Найдено {len(retrieved_docs)} уникальных чанков-кандидатов.")

        if not retrieved_docs:
            return {"answer": "Не найдено релевантных документов.", "status": "not_found", "sources": []}

        final_docs = sorted(retrieved_docs, key=lambda x: x.metadata.get('score', 0), reverse=True)[
                     :self.config["RETRIEVER_TOP_K"]]
        context = ""
        for i, doc in enumerate(final_docs):
            source_info = f"[ИСТОЧНИК {i + 1}: {doc.metadata.get('filename', 'N/A')}, стр. {doc.metadata.get('page', 'N/A')}, слайд {doc.metadata.get('slide', 'N/A')}]"
            context += f"{source_info}\n{doc.page_content}\n\n"

        # --- Новый промпт, запрашивающий JSON ---
        final_prompt = f"""Ты — ассистент-аналитик, который возвращает ответы строго в формате JSON. Проанализируй предоставленный КОНТЕКСТ и ответь на ВОПРОС НА РУССКОМ.

КОНТЕКСТ:
---
{context.strip()}
---

ВОПРОС: "{criterion}"

Твой ответ должен быть ТОЛЬКО JSON объектом со следующей структурой:
{{
  "status": "ОДИН ИЗ СТАТУСОВ: confirmed, not_found, partial, indirect, requires_confirmation",
  "answer": "Твой развернутый ответ на основе контекста, со ссылками на источники в формате [ИСТОЧНИК N] НА РУССКОМ"
}}
"""
        raw_json_string = await self._call_local_llm([{"role": "user", "content": final_prompt}])

        try:
            data = json.loads(raw_json_string)
            clean_answer = data.get("answer", "Ключ 'answer' не найден в ответе модели.")
            status = data.get("status", "requires_confirmation")
            # Простая валидация статуса
            valid_statuses = {"confirmed", "not_found", "partial", "indirect", "requires_confirmation"}
            if status not in valid_statuses:
                status = "requires_confirmation"
        except (json.JSONDecodeError, TypeError):
            # Если модель вернула невалидный JSON или вообще не JSON
            clean_answer = "Ошибка: Модель вернула невалидный JSON. Ответ: " + str(raw_json_string)
            status = "requires_confirmation"

        sources = [
            {"filename": d.metadata.get("filename"), "page": d.metadata.get("page"), "slide": d.metadata.get("slide"),
             "snippet": d.page_content} for d in final_docs]

        return {"answer": clean_answer, "status": status, "sources": sources}


def parse_checklist_from_csv(filename):
    try:
        if not os.path.exists(filename):
            print(f"     Файл чек-листа не найден: {filename}")
            return []
        df = pd.read_csv(filename)
        if 'criterion' not in df.columns:
            print(f"     В файле '{filename}' отсутствует колонка 'criterion'.")
            return []
        criteria = df['criterion'].dropna().astype(str).tolist()
        print(f"    Найдено {len(criteria)} критериев в CSV файле '{filename}'.")
        return criteria
    except Exception as e:
        print(f"    Ошибка при чтении CSV файла '{filename}': {e}")
        return []


async def main(project_names):
    for project_name in project_names:
        print(f"\n======== НАЧИНАЕМ ОБРАБОТКУ ПРОЕКТА: {project_name} ========")

        project_folder = os.path.join(".", project_name)
        documents_path = os.path.join(project_folder, "documents")
        checklist_file = os.path.join(project_folder, f"checklist_{project_name.lower()}.csv")
        index_file_path = os.path.join(project_folder, "vector_index.faiss")
        chunks_file_path = os.path.join(project_folder, "chunks_meta.json")
        report_path = os.path.join(project_folder,
                                   f"verification_report_{project_name}_RAG_FINAL.pdf")

        current_project_config = GLOBAL_CONFIG.copy()
        current_project_config.update({
            "PROJECT_NAME": project_name,
            "PROJECT_FOLDER": project_folder,
            "DOCUMENTS_PATH": documents_path,
            "CHECKLIST_FILE": checklist_file,
            "INDEX_FILE_PATH": index_file_path,
            "CHUNKS_FILE_PATH": chunks_file_path,
        })

        if not os.path.exists(current_project_config["PROJECT_FOLDER"]):
            os.makedirs(current_project_config["PROJECT_FOLDER"])
            os.makedirs(current_project_config["DOCUMENTS_PATH"])
            print(f"Создана структура папок для проекта: {project_name}")

        try:
            system = ComprehensiveRAGSystem(current_project_config)

            print(
                f"\n--- ЭТАП 2: Верификация по чек-листу '{current_project_config['CHECKLIST_FILE']}' для проекта '{project_name}' ---")
            criteria_to_check = parse_checklist_from_csv(current_project_config['CHECKLIST_FILE'])
            if not criteria_to_check:
                print(f"    Пропускаем проект {project_name}: Нет критериев для проверки.")
                continue

            semaphore = asyncio.Semaphore(current_project_config["MAX_CONCURRENT_REQUESTS"])

            async def process_with_semaphore(criterion):
                async with semaphore:
                    result = await system.process_criterion(criterion)
                    await asyncio.sleep(current_project_config["REQUEST_DELAY_SECONDS"])
                    return result

            tasks = [process_with_semaphore(c) for c in criteria_to_check]
            results = await async_tqdm.gather(*tasks, desc=f"Проверка критериев для {project_name}")
            final_report = {c: r for c, r in zip(criteria_to_check, results)}

            print(f"\n--- ИТОГОВЫЙ ОТЧЕТ ВЕРИФИКАЦИИ для проекта {project_name} ---")

            make_pdf(final_report, Path(report_path))

            print(f"\n Отчет сохранен в {report_path}")

        except FileNotFoundError as e:
            print(f"     Ошибка при обработке проекта {project_name}: {e}")
        except Exception as e:
            print(f"     Непредвиденная ошибка при обработке проекта {project_name}: {type(e).__name__}: {e}")

        print(f"\n======== ЗАВЕРШЕНА ОБРАБОТКА ПРОЕКТА: {project_name} ========\n")


if __name__ == "__main__":
    projects_to_process = ["Project_Alfa"] # Оставил один для примера
    asyncio.run(main(projects_to_process))